from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

WHITE = RGBColor(0xFF,0xFF,0xFF)
TITLE_C = RGBColor(0x1B,0x3A,0x5C)
TEXT_C = RGBColor(0x33,0x33,0x33)
ACCENT = RGBColor(0x2E,0x74,0xB5)
SUBTITLE_C = RGBColor(0x66,0x66,0x66)

def bg(slide, color=WHITE):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def txt(slide, l, t, w, h, text, sz=18, bold=False, color=TEXT_C, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align

def bul(slide, l, t, w, h, items, sz=17):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = '\u2022  ' + item
        p.font.size = Pt(sz)
        p.font.color.rgb = TEXT_C
        p.space_before = Pt(10)
        p.space_after = Pt(6)

def bar(slide):
    sh = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(0.06))
    sh.fill.solid()
    sh.fill.fore_color.rgb = ACCENT
    sh.line.fill.background()

def content_slide(title, items):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    bar(s)
    txt(s, 0.8, 0.3, 11.7, 0.8, title, 28, True, TITLE_C)
    bul(s, 0.8, 1.3, 11.7, 5.5, items, 18)

def section_slide(title, subtitle):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, ACCENT)
    txt(s, 1.5, 2.5, 10.3, 1.5, title, 36, True, WHITE, PP_ALIGN.CENTER)
    txt(s, 1.5, 4.0, 10.3, 1.0, subtitle, 20, False, RGBColor(0xCC,0xDD,0xEE), PP_ALIGN.CENTER)

def img_text_slide(title, img_path, items):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    bar(s)
    txt(s, 0.8, 0.3, 11.7, 0.8, title, 28, True, TITLE_C)
    if os.path.exists(img_path):
        s.shapes.add_picture(img_path, Inches(0.8), Inches(1.4), Inches(6.0), Inches(5.0))
    bul(s, 7.2, 1.4, 5.3, 5.5, items, 17)

# S1: Title
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, ACCENT)
txt(s, 1.5, 1.8, 10.3, 1.5, 'Nested Learning:', 44, True, WHITE, PP_ALIGN.CENTER)
txt(s, 1.5, 3.2, 10.3, 1.2, 'The Illusion of Deep Learning Architectures', 34, False, RGBColor(0xDD,0xE8,0xF0), PP_ALIGN.CENTER)
txt(s, 1.5, 5.0, 10.3, 0.6, 'Behrouz, Razaviyayn, Zhong, Mirrokni', 18, False, RGBColor(0xCC,0xDD,0xEE), PP_ALIGN.CENTER)
txt(s, 1.5, 5.6, 10.3, 0.5, 'NeurIPS 2025', 16, False, RGBColor(0xAA,0xCC,0xDD), PP_ALIGN.CENTER)

# S2: Background
content_slide('Research Background & Motivation', [
    'Deep stacking has 4 limits: computational depth unchanged, marginal capacity gain, suboptimal optimization, no continual learning improvement',
    'LLMs after pre-training resemble anterograde amnesia patients - cannot form new long-term memories',
    'Brain achieves continual learning via multi-timescale memory consolidation (online + offline)',
    'Key question: Do we need a new learning paradigm beyond layer stacking?'
])

# S3: LLM Static Nature
content_slide('Core Problem: The Static Nature of LLMs', [
    'LLM knowledge limited to: immediate context window (short-term) or pre-trained MLP params (long-term)',
    'Cannot acquire new knowledge after deployment unless info fits in context window',
    'Existing solutions are expensive, lack generalization, suffer catastrophic forgetting',
    'In-context learning is the only adaptable component, but limited to context window'
])

# S4: Section - NL
section_slide('Nested Learning Paradigm', 'Reinterpreting deep learning as nested optimization problems')

# S5: Brain Inspiration
img_text_slide(
    'Biological Inspiration: Multi-Timescale Brain Updates',
    'extracted/pdf_1770730077/images/b5c6888dccffcf9c49fe36f2e8cc609ffbd6ddb218603aa4457eb7373ee35b68.jpg',
    ['Brain coordinates via Delta~Gamma oscillation bands',
     'Lower neurons: high-freq update; Higher neurons: slow integration',
     'NL reinterprets Transformers as linear layers with different update frequencies',
     'Uniform & reusable structure + multi-timescale = continual learning']
)

# S6: Core Insight
content_slide('Core Insight: Learning = Compressing Context Flow', [
    'All DL components (networks + optimizers) are associative memory systems compressing their own context flow',
    'Associative Memory: operator mapping keys to values, learned by optimizing an objective',
    'Training = acquiring memory that maps data to Local Surprise Signal (LSS)',
    'Components ordered by update frequency form nested multi-level optimization problems'
])

# S7: NL vs DL
img_text_slide(
    'NL Paradigm vs Deep Learning Perspective',
    'extracted/pdf_1770730077/images/109cf9ff8de6674b17b2752024ccc6ac082dfd362a7a626c4e8ae84ea14b5bbf.jpg',
    ['Deep Learning: flattened view, single gradient flow, opaque internals',
     'NL: each sub-component has independent gradient flow & objective',
     'Neural Learning Module: multi-level, each level compresses context at different timescale']
)

# S8: Deep Optimizers
content_slide('Contribution 1: Deep Optimizers', [
    'Momentum is a key-less associative memory compressing past gradients',
    'Adam (with minor modification) is the optimal associative memory for gradients',
    'Ext 1 - Stronger Association: preconditioning (Hessian info)',
    'Ext 2 - Stronger Objective: delta-rule for better capacity management',
    'Ext 3 - Stronger Memory: MLP replaces linear (Deep Momentum GD)',
    'Ext 4 - Nonlinear Output: Newton-Schulz -> Muon optimizer'
])

# S9: Section - HoPE
section_slide('HoPE Architecture', 'Self-Modifying Titans + Continuum Memory System')

# S10: CMS
content_slide('Contributions 2 & 3: Self-Modifying + Continuum Memory', [
    'Self-Modifying Titans: learns to modify its own update algorithm, dynamically adjusts q/k/v',
    'Continuum Memory System (CMS): chain of MLP blocks, each updated at different frequency',
    'CMS generalizes long/short-term memory dichotomy into continuous spectrum',
    'Each MLP params compress own context - representative of abstract knowledge'
])

# S11: HoPE vs Transformer
img_text_slide(
    'HoPE vs Transformer Architecture',
    'extracted/pdf_1770730077/images/7b8edee6945be0eb0adb154a3e8a8c4dea752df596ac0ac3a71edbedd0008043.jpg',
    ['Transformer: only 2 extreme frequency levels',
     'HoPE fills the gap: High/Mid/Low Frequency FFN',
     'Chunk Length: 16 -> 1M -> 16M',
     'Frequency: 1M -> 16 -> 1',
     'Smooth spectrum vs two extremes']
)

# S12: Experiments (two-column)
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
bar(s)
txt(s, 0.8, 0.3, 11.7, 0.8, 'Experimental Results', 28, True, TITLE_C)
txt(s, 0.8, 1.3, 5.5, 0.5, '760M params / 30B tokens', 20, True, ACCENT)
bul(s, 0.8, 1.9, 5.5, 3.0, [
    'HoPE Avg: 52.26% (best)',
    'Titans LMM: 51.56%',
    'Samba: 51.08%',
    'Transformer++: 48.69%'
], 17)
txt(s, 7.0, 1.3, 5.5, 0.5, '1.3B params / 100B tokens', 20, True, ACCENT)
bul(s, 7.0, 1.9, 5.5, 3.0, [
    'HoPE Avg: 57.23% (best)',
    'Titans LMM: 56.82%',
    'Samba: 54.00%',
    'Transformer++: 52.25%'
], 17)
txt(s, 0.8, 5.5, 11.7, 1.0,
    'Key: Dynamic q/k/v projections + deep memory -> lower perplexity & higher accuracy. Advantage grows with scale.',
    16, False, SUBTITLE_C)

# S13: Summary
content_slide('Summary & Critical Evaluation', [
    'NL: unified, math white-box, neuroscience-grounded framework for understanding DL',
    'HoPE consistently outperforms Transformer++ and modern RNNs across scales',
    'Limitation: main paper incomplete due to page limit; many experiments only in appendix',
    'Performance gap modest (57.23 vs 56.82 at 1.3B) - needs larger-scale validation',
    'Open: Can NL guide new architecture design? CMS frequency selection theory?'
])

prs.save('docs/ppt/nested-learning.pptx')
print('Done! 13 slides saved.')
