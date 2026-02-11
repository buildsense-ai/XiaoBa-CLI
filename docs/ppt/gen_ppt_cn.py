from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

WHITE = RGBColor(0xFF,0xFF,0xFF)
TITLE_C = RGBColor(0x1B,0x3A,0x5C)
TEXT_C = RGBColor(0x33,0x33,0x33)
ACCENT = RGBColor(0x2E,0x74,0xB5)
SUBTITLE_C = RGBColor(0x66,0x66,0x66)

def bg(slide, color=WHITE):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def txt(slide, l, t, w, h, text, sz=18, bold=False, color=TEXT_C, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align

def bul(slide, l, t, w, h, items, sz=17):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = '\u2022  ' + item
        p.font.size = Pt(sz)
        p.font.color.rgb = TEXT_C
        p.space_before = Pt(10)
        p.space_after = Pt(6)

def bar(slide):
    sh = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(0.06))
    sh.fill.solid()
    sh.fill.fore_color.rgb = ACCENT
    sh.line.fill.background()

def content_slide(title, items):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    bar(s)
    txt(s, 0.8, 0.3, 11.7, 0.8, title, 28, True, TITLE_C)
    bul(s, 0.8, 1.3, 11.7, 5.5, items, 18)

def section_slide(title, subtitle):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, ACCENT)
    txt(s, 1.5, 2.5, 10.3, 1.5, title, 36, True, WHITE, PP_ALIGN.CENTER)
    txt(s, 1.5, 4.0, 10.3, 1.0, subtitle, 20, False, RGBColor(0xCC,0xDD,0xEE), PP_ALIGN.CENTER)

def img_text_slide(title, img_path, items):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    bar(s)
    txt(s, 0.8, 0.3, 11.7, 0.8, title, 28, True, TITLE_C)
    if os.path.exists(img_path):
        s.shapes.add_picture(img_path, Inches(0.8), Inches(1.4), Inches(6.0), Inches(5.0))
    bul(s, 7.2, 1.4, 5.3, 5.5, items, 17)

# ===== S1: 封面 =====
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, ACCENT)
txt(s, 1.5, 1.8, 10.3, 1.5, 'Nested Learning:', 44, True, WHITE, PP_ALIGN.CENTER)
txt(s, 1.5, 3.2, 10.3, 1.2, '\u6df1\u5ea6\u5b66\u4e60\u67b6\u6784\u7684\u5e7b\u89c9', 34, False, RGBColor(0xDD,0xE8,0xF0), PP_ALIGN.CENTER)
txt(s, 1.5, 5.0, 10.3, 0.6, 'Behrouz, Razaviyayn, Zhong, Mirrokni', 18, False, RGBColor(0xCC,0xDD,0xEE), PP_ALIGN.CENTER)
txt(s, 1.5, 5.6, 10.3, 0.5, 'NeurIPS 2025', 16, False, RGBColor(0xAA,0xCC,0xDD), PP_ALIGN.CENTER)

# ===== S2: 研究背景 =====
content_slide('\u7814\u7a76\u80cc\u666f\u4e0e\u52a8\u673a', [
    '\u6df1\u5ea6\u5b66\u4e60\u5806\u53e0\u5c42\u6570\u7684\u56db\u5927\u5c40\u9650\uff1a\u8ba1\u7b97\u6df1\u5ea6\u4e0d\u53d8\u3001\u5bb9\u91cf\u8fb9\u9645\u9012\u51cf\u3001\u4f18\u5316\u6b21\u4f18\u3001\u6301\u7eed\u5b66\u4e60\u80fd\u529b\u4e0d\u53d8',
    '\u5927\u8bed\u8a00\u6a21\u578b\u9884\u8bad\u7ec3\u540e\u5982\u540c\u201c\u987a\u884c\u6027\u9057\u5fd8\u75c7\u201d\u60a3\u8005\uff0c\u65e0\u6cd5\u5f62\u6210\u65b0\u7684\u957f\u671f\u8bb0\u5fc6',
    '\u5927\u8111\u901a\u8fc7\u591a\u65f6\u95f4\u5c3a\u5ea6\u7684\u8bb0\u5fc6\u5de9\u56fa\u673a\u5236\uff08\u5728\u7ebf\u5de9\u56fa + \u79bb\u7ebf\u5de9\u56fa\uff09\u5b9e\u73b0\u6301\u7eed\u5b66\u4e60',
    '\u6838\u5fc3\u8ffd\u95ee\uff1a\u662f\u5426\u9700\u8981\u8d85\u8d8a\u201c\u5c42\u5806\u53e0\u201d\u7684\u65b0\u5b66\u4e60\u8303\u5f0f\uff1f',
])

# ===== S3: LLM\u9759\u6001\u672c\u8d28 =====
content_slide('\u6838\u5fc3\u95ee\u9898\uff1a\u5927\u8bed\u8a00\u6a21\u578b\u7684\u9759\u6001\u672c\u8d28', [
    '\u77e5\u8bc6\u5c40\u9650\u4e8e\u5373\u65f6\u4e0a\u4e0b\u6587\u7a97\u53e3\uff08\u77ed\u671f\uff09\u6216\u9884\u8bad\u7ec3 MLP \u53c2\u6570\uff08\u957f\u671f\uff09',
    '\u90e8\u7f72\u540e\u65e0\u6cd5\u6301\u7eed\u83b7\u53d6\u65b0\u77e5\u8bc6\uff0c\u9664\u975e\u4fe1\u606f\u4ecd\u5728\u4e0a\u4e0b\u6587\u7a97\u53e3\u5185',
    '\u73b0\u6709\u65b9\u6848\uff08\u5fae\u8c03\u3001\u5916\u90e8\u8bb0\u5fc6\u7b49\uff09\u8ba1\u7b97\u6602\u8d35\u3001\u7f3a\u4e4f\u6cdb\u5316\u3001\u6613\u707e\u96be\u6027\u9057\u5fd8',
    '\u4e0a\u4e0b\u6587\u5b66\u4e60\u662f\u552f\u4e00\u53ef\u9002\u5e94\u7ec4\u4ef6\uff0c\u4f46\u53d7\u9650\u4e8e\u7a97\u53e3\u5927\u5c0f',
])

# ===== S4: \u7ae0\u8282\u5206\u9694 =====
section_slide('Nested Learning \u8303\u5f0f', '\u5c06\u6df1\u5ea6\u5b66\u4e60\u91cd\u65b0\u7406\u89e3\u4e3a\u5d4c\u5957\u4f18\u5316\u95ee\u9898\u7684\u96c6\u5408')

# ===== S5: \u751f\u7269\u542f\u53d1 =====
img_text_slide(
    '\u751f\u7269\u542f\u53d1\uff1a\u5927\u8111\u7684\u591a\u65f6\u95f4\u5c3a\u5ea6\u66f4\u65b0',
    'extracted/pdf_1770730077/images/b5c6888dccffcf9c49fe36f2e8cc609ffbd6ddb218603aa4457eb7373ee35b68.jpg',
    ['\u5927\u8111\u901a\u8fc7 Delta~Gamma \u4e94\u79cd\u9891\u6bb5\u534f\u8c03\u795e\u7ecf\u6d3b\u52a8',
     '\u4f4e\u5c42\u795e\u7ecf\u5143\u9ad8\u9891\u66f4\u65b0\uff0c\u9ad8\u5c42\u795e\u7ecf\u5143\u4f4e\u9891\u6574\u5408',
     'NL \u5c06 Transformer \u91cd\u65b0\u89e3\u8bfb\u4e3a\u6309\u4e0d\u540c\u9891\u7387\u66f4\u65b0\u7684\u7ebf\u6027\u5c42',
     '\u7edf\u4e00\u53ef\u590d\u7528\u7ed3\u6784 + \u591a\u65f6\u95f4\u5c3a\u5ea6 = \u6301\u7eed\u5b66\u4e60']
)

# ===== S6: \u6838\u5fc3\u6d1e\u5bdf =====
content_slide('\u6838\u5fc3\u6d1e\u5bdf\uff1a\u5b66\u4e60 = \u538b\u7f29\u4e0a\u4e0b\u6587\u6d41', [
    '\u6240\u6709\u6df1\u5ea6\u5b66\u4e60\u7ec4\u4ef6\uff08\u7f51\u7edc + \u4f18\u5316\u5668\uff09\u90fd\u662f\u538b\u7f29\u81ea\u8eab\u4e0a\u4e0b\u6587\u6d41\u7684\u5173\u8054\u8bb0\u5fc6\u7cfb\u7edf',
    '\u5173\u8054\u8bb0\u5fc6\uff1a\u5c06\u952e\uff08keys\uff09\u6620\u5c04\u5230\u503c\uff08values\uff09\u7684\u7b97\u5b50\uff0c\u901a\u8fc7\u4f18\u5316\u76ee\u6807\u51fd\u6570\u5b66\u4e60\u6620\u5c04',
    '\u8bad\u7ec3 = \u5c06\u6570\u636e\u6620\u5c04\u5230\u201c\u5c40\u90e8\u60ca\u8bb6\u4fe1\u53f7\u201d\uff08LSS\uff09\u7684\u8bb0\u5fc6\u83b7\u53d6\u8fc7\u7a0b',
    '\u4e0d\u540c\u7ec4\u4ef6\u6309\u66f4\u65b0\u9891\u7387\u6392\u5e8f\uff0c\u5f62\u6210\u5d4c\u5957\u7684\u591a\u5c42\u7ea7\u4f18\u5316\u95ee\u9898',
])

# ===== S7: NL vs DL =====
img_text_slide(
    'NL \u8303\u5f0f vs \u6df1\u5ea6\u5b66\u4e60\u89c6\u89d2',
    'extracted/pdf_1770730077/images/109cf9ff8de6674b17b2752024ccc6ac082dfd362a7a626c4e8ae84ea14b5bbf.jpg',
    ['\u6df1\u5ea6\u5b66\u4e60\uff1a\u6241\u5e73\u5316\u89c6\u89d2\uff0c\u5355\u4e00\u68af\u5ea6\u6d41\uff0c\u5185\u90e8\u52a8\u6001\u4e0d\u900f\u660e',
     'NL\uff1a\u6bcf\u4e2a\u5b50\u7ec4\u4ef6\u62e5\u6709\u72ec\u7acb\u68af\u5ea6\u6d41\u548c\u4f18\u5316\u76ee\u6807\uff0c\u6570\u5b66\u767d\u76d2',
     '\u795e\u7ecf\u5b66\u4e60\u6a21\u5757\uff1a\u591a\u5c42\u7ea7\u7ed3\u6784\uff0c\u6bcf\u5c42\u538b\u7f29\u4e0d\u540c\u65f6\u95f4\u5c3a\u5ea6\u7684\u4e0a\u4e0b\u6587']
)

# ===== S8: Deep Optimizers =====
content_slide('\u8d21\u732e\u4e00\uff1a\u6df1\u5ea6\u4f18\u5316\u5668\uff08Deep Optimizers\uff09', [
    '\u52a8\u91cf\u9879\u672c\u8d28\u4e0a\u662f\u538b\u7f29\u8fc7\u53bb\u68af\u5ea6\u7684\u65e0\u952e\u5173\u8054\u8bb0\u5fc6\u6a21\u5757',
    'Adam\uff08\u7ecf\u5c0f\u4fee\u6539\uff09\u662f\u68af\u5ea6\u7684\u6700\u4f18\u5173\u8054\u8bb0\u5fc6',
    '\u6269\u5c55\u4e00\uff1a\u66f4\u5f3a\u5173\u8054\u2014\u2014\u9884\u8c03\u8282\uff08\u5229\u7528 Hessian \u4fe1\u606f\uff09',
    '\u6269\u5c55\u4e8c\uff1a\u66f4\u5f3a\u76ee\u6807\u2014\u2014delta-rule\uff0c\u66f4\u597d\u7684\u5bb9\u91cf\u7ba1\u7406',
    '\u6269\u5c55\u4e09\uff1a\u66f4\u5f3a\u8bb0\u5fc6\u2014\u2014MLP \u66ff\u4ee3\u7ebf\u6027\u5c42\uff08\u6df1\u5ea6\u52a8\u91cf\u68af\u5ea6\u4e0b\u964d\uff09',
    '\u6269\u5c55\u56db\uff1a\u975e\u7ebf\u6027\u8f93\u51fa\u2014\u2014Newton-Schulz \u65b9\u6cd5\uff0c\u7b49\u4ef7\u4e8e Muon \u4f18\u5316\u5668',
])

# ===== S9: \u7ae0\u8282\u5206\u9694 =====
section_slide('HoPE \u67b6\u6784', '\u81ea\u4fee\u6539 Titans + \u8fde\u7eed\u8bb0\u5fc6\u7cfb\u7edf')

# ===== S10: CMS =====
content_slide('\u8d21\u732e\u4e8c\u4e09\uff1a\u81ea\u4fee\u6539\u6a21\u578b + \u8fde\u7eed\u8bb0\u5fc6\u7cfb\u7edf', [
    '\u81ea\u4fee\u6539 Titans\uff1a\u5b66\u4e60\u4fee\u6539\u81ea\u8eab\u66f4\u65b0\u7b97\u6cd5\uff0c\u52a8\u6001\u8c03\u6574 q/k/v \u6295\u5f71',
    '\u8fde\u7eed\u8bb0\u5fc6\u7cfb\u7edf\uff08CMS\uff09\uff1a\u4e00\u6761 MLP \u94fe\uff0c\u6bcf\u4e2a\u5757\u4ee5\u4e0d\u540c\u9891\u7387\u66f4\u65b0',
    'CMS \u5c06\u4f20\u7edf\u957f\u77ed\u671f\u8bb0\u5fc6\u7684\u4e8c\u5206\u6cd5\u6cdb\u5316\u4e3a\u8fde\u7eed\u9891\u8c31',
    '\u6bcf\u5c42 MLP \u53c2\u6570\u538b\u7f29\u81ea\u8eab\u4e0a\u4e0b\u6587\uff0c\u662f\u8be5\u65f6\u95f4\u5c3a\u5ea6\u4e0a\u62bd\u8c61\u77e5\u8bc6\u7684\u4ee3\u8868',
])

# ===== S11: HoPE vs Transformer =====
img_text_slide(
    'HoPE vs Transformer \u67b6\u6784\u5bf9\u6bd4',
    'extracted/pdf_1770730077/images/7b8edee6945be0eb0adb154a3e8a8c4dea752df596ac0ac3a71edbedd0008043.jpg',
    ['Transformer \u53ea\u6709\u4e24\u4e2a\u6781\u7aef\u9891\u7387\u5c42\u7ea7\uff08FFN \u53c2\u6570\u9891\u7387=0\uff0c\u6ce8\u610f\u529b\u6fc0\u6d3b\u9891\u7387=\u221e\uff09',
     'HoPE \u586b\u5145\u4e2d\u95f4\u9891\u8c31\uff1a\u9ad8/\u4e2d/\u4f4e\u9891 FFN \u5c42',
     '\u5757\u957f\u5ea6\uff1a16 \u2192 1M \u2192 16M\uff08\u4ece\u5e95\u5230\u9876\uff09',
     '\u66f4\u65b0\u9891\u7387\uff1a1M \u2192 16 \u2192 1\uff08\u4ece\u5feb\u5230\u6162\uff09',
     '\u5e73\u6ed1\u7684\u8bb0\u5fc6\u65f6\u95f4\u5c3a\u5ea6\u9891\u8c31 vs \u4e24\u6781\u5206\u5316']
)

# ===== S12: \u5b9e\u9a8c\u7ed3\u679c =====
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
bar(s)
txt(s, 0.8, 0.3, 11.7, 0.8, '\u5b9e\u9a8c\u7ed3\u679c\uff1a\u8bed\u8a00\u5efa\u6a21\u4e0e\u5e38\u8bc6\u63a8\u7406', 28, True, TITLE_C)
txt(s, 0.8, 1.3, 5.5, 0.5, '760M \u53c2\u6570 / 30B tokens', 20, True, ACCENT)
bul(s, 0.8, 1.9, 5.5, 3.0, [
    'HoPE \u5e73\u5747\uff1a52.26%\uff08\u6700\u4f18\uff09',
    'Titans LMM\uff1a51.56%',
    'Samba\uff1a51.08%',
    'Transformer++\uff1a48.69%',
], 17)
txt(s, 7.0, 1.3, 5.5, 0.5, '1.3B \u53c2\u6570 / 100B tokens', 20, True, ACCENT)
bul(s, 7.0, 1.9, 5.5, 3.0, [
    'HoPE \u5e73\u5747\uff1a57.23%\uff08\u6700\u4f18\uff09',
    'Titans LMM\uff1a56.82%',
    'Samba\uff1a54.00%',
    'Transformer++\uff1a52.25%',
], 17)
txt(s, 0.8, 5.5, 11.7, 1.0,
    '\u5173\u952e\uff1a\u52a8\u6001 q/k/v \u6295\u5f71 + \u6df1\u5ea6\u8bb0\u5fc6\u6a21\u5757 \u2192 \u66f4\u4f4e\u56f0\u60d1\u5ea6\u3001\u66f4\u9ad8\u51c6\u786e\u7387\u3002\u4f18\u52bf\u968f\u89c4\u6a21\u589e\u5927\u800c\u66f4\u660e\u663e\u3002',
    16, False, SUBTITLE_C)

# ===== S13: \u603b\u7ed3 =====
content_slide('\u603b\u7ed3\u4e0e\u8bc4\u4ef7', [
    'NL \u63d0\u4f9b\u4e86\u6570\u5b66\u767d\u76d2\u3001\u795e\u7ecf\u79d1\u5b66\u5408\u7406\u7684\u7edf\u4e00\u6846\u67b6\uff0c\u540c\u65f6\u89e3\u91ca\u67b6\u6784\u8bbe\u8ba1\u548c\u4f18\u5316\u7b97\u6cd5',
    'HoPE \u5728\u591a\u4e2a\u89c4\u6a21\u4e0b\u6301\u7eed\u8d85\u8d8a Transformer++ \u548c\u73b0\u4ee3\u5faa\u73af\u7f51\u7edc',
    '\u5c40\u9650\uff1a\u4e3b\u6587\u56e0\u9875\u6570\u9650\u5236\u4e0d\u591f\u5b8c\u6574\uff0c\u5927\u91cf\u5b9e\u9a8c\u4ec5\u5728\u9644\u5f55\u4e2d',
    '\u6027\u80fd\u5dee\u8ddd\u8f83\u5c0f\uff0857.23 vs 56.82 @ 1.3B\uff09\uff0c\u9700\u8981\u66f4\u5927\u89c4\u6a21\u9a8c\u8bc1',
    '\u5f00\u653e\u95ee\u9898\uff1aNL \u80fd\u5426\u6307\u5bfc\u8bbe\u8ba1\u5168\u65b0\u67b6\u6784\uff1fCMS \u9891\u7387\u9009\u62e9\u662f\u5426\u6709\u7406\u8bba\u6307\u5bfc\uff1f',
])

prs.save('docs/ppt/nested-learning-cn.pptx')
print('Done! 13 slides saved to docs/ppt/nested-learning-cn.pptx')
