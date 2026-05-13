"""
从 PPT 文件中提取文本和图片内容。
支持混合模式：python-pptx 提取文字 + win32com 导出图片 + Vision API 分析。
"""

import argparse
import json
import os
import sys
from pathlib import Path

from pptx import Presentation
from pptx.shapes.base import BaseShape
from pptx.enum.shapes import MSO_SHAPE_TYPE


def extract_text_from_shape(shape: BaseShape) -> str:
    """递归提取 shape 中的所有文字。"""
    texts = []

    if shape.has_text_frame:
        for paragraph in shape.text_frame.paragraphs:
            line = paragraph.text.strip()
            if line:
                texts.append(line)

    if shape.has_table:
        table = shape.table
        for row in table.rows:
            row_texts = []
            for cell in row.cells:
                cell_text = cell.text.strip().replace('\n', ' ')
                if cell_text:
                    row_texts.append(cell_text)
            if row_texts:
                texts.append(' | '.join(row_texts))

    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for child in shape.shapes:
            texts.append(extract_text_from_shape(child))

    return '\n'.join(texts)


def detect_visual(shape: BaseShape) -> bool:
    """判断 shape 是否包含图片或图表。"""
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        return True
    if shape.shape_type == MSO_SHAPE_TYPE.CHART:
        return True
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for child in shape.shapes:
            if detect_visual(child):
                return True
    return False


def extract_pptx_text(pptx_path: str) -> dict:
    """提取 PPT 所有页面的文字，并检测图片/图表。"""
    prs = Presentation(str(pptx_path))
    slides_data = []

    for i, slide in enumerate(prs.slides):
        texts = []
        has_visual = False

        for shape in slide.shapes:
            texts.append(extract_text_from_shape(shape))
            if detect_visual(shape):
                has_visual = True

        slide_text = '\n'.join(t for t in texts if t)
        slides_data.append({
            "index": i + 1,
            "text": slide_text,
            "has_visual": has_visual,
            "image_path": None,
            "vision_analysis": None,
        })

    return {
        "file_name": os.path.basename(pptx_path),
        "slide_count": len(slides_data),
        "slides": slides_data,
    }


def export_slides_as_images(pptx_path: str, slides_data: list, output_dir: str) -> None:
    """用 win32com 将有图片的页面导出为 PNG。"""
    import pythoncom
    pythoncom.CoInitialize()

    from win32com.client import Dispatch

    ppt_abs = str(Path(pptx_path).resolve())
    out_abs = str(Path(output_dir).resolve())
    os.makedirs(out_abs, exist_ok=True)

    powerpoint = None
    presentation = None
    errors = []

    try:
        powerpoint = Dispatch("PowerPoint.Application")
        powerpoint.Visible = True
        presentation = powerpoint.Presentations.Open(ppt_abs, WithWindow=False)

        visual_indices = [s["index"] for s in slides_data if s["has_visual"]]
        if not visual_indices:
            return

        for idx in visual_indices:
            try:
                image_path = os.path.join(out_abs, f"slide_{idx}.png")
                # 如果已经存在就跳过
                if os.path.exists(image_path):
                    slides_data[idx - 1]["image_path"] = image_path
                    continue

                slide = presentation.Slides(idx)
                slide.Export(image_path, "PNG", 1920, 1080)
                slides_data[idx - 1]["image_path"] = image_path
            except Exception as e:
                errors.append(f"Slide {idx}: {e}")

    finally:
        if presentation:
            try:
                presentation.Close()
            except Exception:
                pass
        if powerpoint:
            try:
                powerpoint.Quit()
            except Exception:
                pass
        try:
            pythoncom.CoUninitialize()
        except Exception:
            pass

    if errors:
        print(f"[警告] 部分页面导出失败: {'; '.join(errors)}", file=sys.stderr)


def call_vision_service(image_path: str, prompt: str, base_url: str):  # -> str | None:
    """调用 advanced-reader-service 分析图片。"""
    import requests

    try:
        with open(image_path, 'rb') as f:
            response = requests.post(
                f"{base_url}/analyze",
                files={"file": (os.path.basename(image_path), f, 'image/png')},
                data={"prompt": prompt},
                timeout=120,
            )
        if response.status_code == 200:
            data = response.json()
            return data.get("analysis", "")
        else:
            print(f"[警告] Vision API 返回 {response.status_code}: {response.text[:200]}", file=sys.stderr)
            return None
    except Exception as e:
        print(f"[警告] Vision API 调用失败: {e}", file=sys.stderr)
        return None


def main():
    parser = argparse.ArgumentParser(description="提取 PPT 文本和图片内容")
    parser.add_argument("pptx_path", help="PPT 文件路径")
    parser.add_argument("-o", "--output", default=None, help="输出 JSON 文件路径（默认输出到 stdout）")
    parser.add_argument("--with-vision", action="store_true", help="对包含图片的页面调用 Vision API 分析")
    parser.add_argument("--vision-url", default="http://localhost:8000", help="Advanced Reader Service URL")
    parser.add_argument("--vision-prompt", default="请详细描述这张幻灯片中的图片、图表、表格的内容。如果包含数据，请提取出所有数据。如果包含文字，请逐字识别。", help="Vision API 的 prompt")
    parser.add_argument("--image-dir", default=None, help="导出图片的目录（默认在 PPT 同目录下创建 slides_img 文件夹）")
    args = parser.parse_args()

    pptx_path = args.pptx_path
    if not os.path.exists(pptx_path):
        print(f"错误：文件不存在 - {pptx_path}", file=sys.stderr)
        sys.exit(1)

    if not pptx_path.lower().endswith(('.ppt', '.pptx')):
        print("错误：仅支持 .ppt 或 .pptx 文件", file=sys.stderr)
        sys.exit(1)

    if pptx_path.lower().endswith('.ppt'):
        print("注意：.ppt 格式只能用 win32com 导出图片，无法直接提取文字。", file=sys.stderr)
        print("请将文件另存为 .pptx 格式后重试。", file=sys.stderr)
        sys.exit(1)

    # Step 1: 提取文字
    print("正在提取 PPT 文字...", file=sys.stderr)
    result = extract_pptx_text(pptx_path)
    visual_count = sum(1 for s in result["slides"] if s["has_visual"])
    print(f"共 {result['slide_count']} 页，其中 {visual_count} 页包含图片/图表", file=sys.stderr)

    # Step 2: 对包含图片的页面导出 PNG
    if args.with_vision and visual_count > 0:
        image_dir = args.image_dir or os.path.join(os.path.dirname(pptx_path) or ".", "slides_img")
        print(f"正在导出图片页面...", file=sys.stderr)
        export_slides_as_images(pptx_path, result["slides"], image_dir)

        # Step 3: 调用 Vision API
        exported = [s for s in result["slides"] if s["image_path"]]
        print(f"成功导出 {len(exported)} 张图片", file=sys.stderr)

        for slide in exported:
            print(f"正在分析第 {slide['index']} 页...", file=sys.stderr)
            analysis = call_vision_service(slide["image_path"], args.vision_prompt, args.vision_url)
            if analysis:
                slide["vision_analysis"] = analysis

    # 输出
    json_output = json.dumps(result, ensure_ascii=False, indent=2)
    if args.output:
        with open(args.output, 'w', encoding='utf-8') as f:
            f.write(json_output)
        print(f"结果已保存到 {args.output}", file=sys.stderr)
    else:
        print(json_output)


if __name__ == "__main__":
    main()
